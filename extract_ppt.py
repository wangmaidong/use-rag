from pptx import Presentation

def extract_ppt_text(file_path):
    """
    提取PPT文件中的所有文本内容，并以字符串返回。
    :param file_path: PPT文件路径
    :return: 所有文本内容（以换行符分隔）
    """
    ppt = Presentation(file_path)
    text_list = []
    for slide in ppt.slides:
        for shape in slide.shapes:
            if hasattr(shape, 'text'):
                text_list.append(shape.text)
    all_text = '\n'.join(text_list)
    return all_text

if __name__ == "__main__":
    ppt_file = "example_file/example.pptx"
    result = extract_ppt_text(ppt_file)
    print(result)