# 导入PyMuPDF库（fitz），用于处理PDF文件
import fitz  # PyMuPDF
# 导入Optional类型提示
from typing import Optional
# 导入日志logging功能
import logging

# 获取当前模块日志记录器
logger = logging.getLogger(__name__)

# 定义用于提取PDF所有文本内容的函数
def extract_pdf_text(pdf_path: str) -> str:
    """
    提取PDF文件中的所有文本内容

    参数:
        pdf_path (str): PDF文件路径

    返回:
        str: 合并后的所有页文本

    异常:
        FileNotFoundError: 文件不存在
        Exception: PDF文件读取失败
    """
    try:
        # 打开PDF文件
        pdf = fitz.open(pdf_path)
        try:
            # 新建一个空列表，用来存储每页文本
            text_list = []
            # 遍历每一页
            for page in pdf:
                # 获取当前页文本，并加入列表
                text_list.append(page.get_text("text"))  # type: ignore
            # 将每页文本用换行拼接成一个大字符串
            all_text = "\n".join(text_list)
            # 返回拼接后的文本
            return all_text
        finally:
            # 确保关闭PDF文件
            pdf.close()
    except FileNotFoundError:
        # 如果文件未找到，记录错误日志
        logger.error(f"PDF文件不存在: {pdf_path}")
        # 向上抛出异常
        raise
    except Exception as e:
        # 其他异常情况，记录错误信息
        logger.error(f"提取PDF文本失败: {pdf_path}, 错误: {str(e)}")
        # 抛出异常
        raise

# 导入python-docx的Document类
from docx import Document

# 定义提取Word文档所有段落文本的函数
def extract_text_from_word(file_path: str) -> str:
    """
    从Word文档中提取所有段落的文本，并以字符串返回。

    参数:
        file_path (str): Word文档的路径

    返回:
        str: 文本内容字符串

    异常:
        FileNotFoundError: 文件不存在
        Exception: Word文件读取失败
    """
    try:
        # 加载Word文档
        doc = Document(file_path)
        # 取所有段落的文本，并用换行符拼接
        text = "\n".join([para.text for para in doc.paragraphs])
        # 返回拼接好的文本
        return text
    except FileNotFoundError:
        # 文件未找到时记录日志
        logger.error(f"Word文件不存在: {file_path}")
        # 抛出异常
        raise
    except Exception as e:
        # 其它异常记录错误信息
        logger.error(f"提取Word文本失败: {file_path}, 错误: {str(e)}")
        # 抛出异常
        raise

# 导入openpyxl库，用于操作Excel文件
import openpyxl

# 定义函数提取Excel文件中的所有文本
def extract_text_from_excel(file_path: str) -> str:
    """
    从Excel文件中提取所有单元格内容为文本，并以字符串返回。

    参数:
        file_path (str): Excel文件路径

    返回:
        str: 文本内容字符串

    异常:
        FileNotFoundError: 文件不存在
        Exception: Excel文件读取失败
    """
    try:
        # 加载Excel工作簿
        wb = openpyxl.load_workbook(file_path, data_only=True)
        try:
            # 取得活动工作表
            ws = wb.active
            # 新建空列表保存每一行字符串
            rows = []
            # 遍历所有行，只取单元格的值
            for row in ws.iter_rows(values_only=True):
                # 将每行单元格内容用Tab连接，空值转换为空字符串
                rows.append("\t".join([str(cell) if cell is not None else "" for cell in row]))
            # 用换行符拼接所有行
            all_text = "\n".join(rows)
            # 返回最终文本
            return all_text
        finally:
            # 关闭Excel工作簿
            wb.close()
    except FileNotFoundError:
        # 文件未找到时日志记录
        logger.error(f"Excel文件不存在: {file_path}")
        raise
    except Exception as e:
        # 其它异常日志并抛出
        logger.error(f"提取Excel文本失败: {file_path}, 错误: {str(e)}")
        raise

# 导入python-pptx库的Presentation类
from pptx import Presentation

# 定义函数提取PPT文件所有文本内容
def extract_ppt_text(file_path: str) -> str:
    """
    提取PPT文件中的所有文本内容，并以字符串返回。

    参数:
        file_path (str): PPT文件路径

    返回:
        str: 所有文本内容（以换行符分隔）

    异常:
        FileNotFoundError: 文件不存在
        Exception: PPT文件读取失败
    """
    try:
        # 加载PPT文件
        ppt = Presentation(file_path)
        # 新建列表存储所有文本内容
        text_list = []
        # 遍历PPT中的每张幻灯片
        for slide in ppt.slides:
            # 遍历当前幻灯片的每个形状
            for shape in slide.shapes:
                # 判断是否含有文本，且文本不为空
                if hasattr(shape, "text") and shape.text.strip():
                    # 有文本时加入结果列表
                    text_list.append(shape.text)
        # 用换行符拼接所有文本
        all_text = "\n".join(text_list)
        # 返回所有文本内容
        return all_text
    except FileNotFoundError:
        # 文件未找到时日志打印
        logger.error(f"PPT文件不存在: {file_path}")
        raise
    except Exception as e:
        # 处理其它异常
        logger.error(f"提取PPT文本失败: {file_path}, 错误: {str(e)}")
        raise

# 导入BeautifulSoup用于解析HTML
from bs4 import BeautifulSoup  # BeautifulSoup用于解析HTML

# 定义函数，从HTML文件提取所有文本内容
def extract_text_from_html(file_path: str) -> str:
    """
    从指定HTML文件中提取所有文本内容

    参数:
        file_path (str): HTML文件路径

    返回:
        str: 提取的文本内容

    异常:
        FileNotFoundError: 文件不存在
        Exception: HTML文件读取失败
    """
    try:
        # 以utf-8编码方式打开HTML文件
        with open(file_path, "r", encoding="utf-8") as f:
            # 读取HTML文件所有内容
            html = f.read()
        # 创建BeautifulSoup对象
        soup = BeautifulSoup(html, "html.parser")
        # 用换行分隔符获取全部文本
        text = soup.get_text(separator="\n", strip=True)
        # 返回文本
        return text
    except FileNotFoundError:
        # 文件不存在，记录日志
        logger.error(f"HTML文件不存在: {file_path}")
        raise
    except Exception as e:
        # 其它异常记录并抛出
        logger.error(f"提取HTML文本失败: {file_path}, 错误: {str(e)}")
        raise

# 导入内置json库
import json

# 定义提取JSON文件文本内容的函数
def extract_text_from_json(filename: str) -> str:
    """
    从JSON文件中提取文本内容并格式化为字符串

    参数:
        filename (str): JSON文件路径

    返回:
        str: 格式化后的JSON文本内容

    异常:
        FileNotFoundError: 文件不存在
        json.JSONDecodeError: JSON解析失败
    """
    try:
        # 以utf-8编码打开JSON文件
        with open(filename, "r", encoding="utf-8") as f:
            # 加载JSON内容到Python对象
            data = json.load(f)
        # 格式化JSON为缩进文本，显示中文
        text = json.dumps(data, ensure_ascii=False, indent=2)
        # 返回字符串格式JSON内容
        return text
    except FileNotFoundError:
        # 文件不存在时记录日志
        logger.error(f"JSON文件不存在: {filename}")
        raise
    except json.JSONDecodeError as e:
        # JSON解析异常日志
        logger.error(f"JSON解析失败: {filename}, 错误: {str(e)}")
        raise

# 导入lxml库的etree模块用于XML处理
from lxml import etree

# 定义函数，从XML文件提取所有文本内容
def extract_xml_text(file_path: str) -> str:
    """
    读取XML文件并提取所有文本内容

    参数:
        file_path (str): XML文件路径

    返回:
        str: 提取的所有文本内容

    异常:
        FileNotFoundError: 文件不存在
        etree.XMLSyntaxError: XML解析失败
    """
    try:
        # 用utf-8编码打开XML文件
        with open(file_path, "r", encoding="utf-8") as f:
            # 读取XML字符串内容
            xml = f.read()
        # 解析为XML树结构对象
        root = etree.fromstring(xml.encode("utf-8"))
        # 遍历所有文本节点并用空格拼接
        text = " ".join(root.itertext())
        # 返回拼接后的文本
        return text
    except FileNotFoundError:
        # 文件不存在日志
        logger.error(f"XML文件不存在: {file_path}")
        raise
    except etree.XMLSyntaxError as e:
        # XML语法异常日志
        logger.error(f"XML解析失败: {file_path}, 错误: {str(e)}")
        raise
    except Exception as e:
        # 其它异常日志
        logger.error(f"提取XML文本失败: {file_path}, 错误: {str(e)}")
        raise

# 导入csv模块
import csv

# 定义读取CSV内容并串成字符串的函数
def read_csv_to_text(filename: str) -> str:
    """
    读取CSV文件内容，并将每行用逗号连接，所有行用换行符拼接成一个字符串返回。

    参数:
        filename (str): CSV文件路径

    返回:
        str: 拼接后的字符串

    异常:
        FileNotFoundError: 文件不存在
    """
    try:
        # 以utf-8编码方式打开CSV文件
        with open(filename, "r", encoding="utf-8") as f:
            # 创建csv.reader对象逐行读取
            reader = csv.reader(f)
            # 每行用逗号拼接并放到列表
            rows = [", ".join(row) for row in reader]
        # 用换行拼接所有行
        all_text = "\n".join(rows)
        # 返回结果
        return all_text
    except FileNotFoundError:
        # 文件不存在日志
        logger.error(f"CSV文件不存在: {filename}")
        raise
    except Exception as e:
        # 其它异常日志
        logger.error(f"读取CSV文件失败: {filename}, 错误: {str(e)}")
        raise

# 定义读取文本文件内容的函数
def read_text_file(filename: str) -> str:
    """
    读取指定文本文件内容并返回

    参数:
        filename (str): 文件路径

    返回:
        str: 文件内容字符串

    异常:
        FileNotFoundError: 文件不存在
    """
    try:
        # 以utf-8只读方式打开文本文件
        with open(filename, "r", encoding="utf-8") as f:
            # 读取文件的所有内容
            text = f.read()
        # 返回字符串
        return text
    except FileNotFoundError:
        # 文件未找到记录日志
        logger.error(f"文本文件不存在: {filename}")
        raise
    except Exception as e:
        # 其它异常情况日志记录
        logger.error(f"读取文本文件失败: {filename}, 错误: {str(e)}")
        raise

# 定义读取Markdown文件内容的函数
def read_markdown_file(file_path: str) -> str:
    """
    读取Markdown文件内容并返回

    参数:
        file_path (str): Markdown文件路径

    返回:
        str: 文件内容字符串

    异常:
        FileNotFoundError: 文件不存在
    """
    try:
        # 以utf-8编码只读打开Markdown文件
        with open(file_path, "r", encoding="utf-8") as f:
            # 读取并返回全部内容
            return f.read()
    except FileNotFoundError:
        # 文件不存在日志
        logger.error(f"Markdown文件不存在: {file_path}")
        raise
    except Exception as e:
        # 其它异常日志
        logger.error(f"读取Markdown文件失败: {file_path}, 错误: {str(e)}")
        raise
